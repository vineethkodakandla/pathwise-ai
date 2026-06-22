"""Generate PathWise AI presentation deck."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Brand colors
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
TEAL = RGBColor(0x00, 0xA8, 0xB5)
ACCENT = RGBColor(0xFF, 0x8C, 0x42)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
GRAY = RGBColor(0x5A, 0x6B, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0xA6, 0x4D)
YELLOW = RGBColor(0xF0, 0xB4, 0x29)
RED = RGBColor(0xD1, 0x3B, 0x3B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]


def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_bar(slide, color=TEAL, height=Inches(0.18)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_footer(slide, page_num):
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.3))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "PathWise AI  |  Team Pathfinders  |  COSC 6370-001"
    r.font.size = Pt(10)
    r.font.color.rgb = GRAY
    # right side page num
    tb2 = slide.shapes.add_textbox(Inches(12), Inches(7.1), Inches(1), Inches(0.3))
    tf2 = tb2.text_frame
    tf2.margin_left = tf2.margin_right = 0
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = str(page_num)
    r2.font.size = Pt(10)
    r2.font.color.rgb = GRAY


def add_title(slide, title, subtitle=None):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.9))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = NAVY
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(15)
        r2.font.color.rgb = TEAL
        r2.font.italic = True
    # underline accent
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.3), Inches(0.8), Inches(0.06))
    acc.fill.solid()
    acc.fill.fore_color.rgb = ACCENT
    acc.line.fill.background()


def add_bullets(slide, items, left=0.6, top=1.6, width=12.1, height=5.2, size=16, spacing=10):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(spacing)
        if isinstance(item, tuple):
            text, bold = item
        else:
            text, bold = item, False
        # bullet
        rb = p.add_run()
        rb.text = "\u25B8  "
        rb.font.size = Pt(size)
        rb.font.color.rgb = TEAL
        rb.font.bold = True
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = NAVY
        r.font.bold = bold


def add_text_block(slide, text, left, top, width, height, size=14, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return tf


def add_table(slide, headers, rows, left, top, width, height, header_color=NAVY, header_text=WHITE, body_color=NAVY, font_size=12):
    cols = len(headers)
    n_rows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(n_rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    tbl = tbl_shape.table
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        tf = cell.text_frame
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        tf.text = ""
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = h
        r.font.size = Pt(font_size + 1)
        r.font.bold = True
        r.font.color.rgb = header_text
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if ri % 2 == 1 else WHITE
            tf = cell.text_frame
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            tf.word_wrap = True
            tf.text = ""
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = val
            r.font.size = Pt(font_size)
            r.font.color.rgb = body_color
    return tbl


def add_card(slide, left, top, width, height, title, body, accent=TEAL):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = accent
    card.line.width = Pt(1.5)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = accent
    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    r2 = p2.add_run()
    r2.text = body
    r2.font.size = Pt(11)
    r2.font.color.rgb = NAVY


# =========================================================================
# SLIDE 1 — TITLE
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
# accent stripe
stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), prs.slide_width, Inches(0.08))
stripe.fill.solid(); stripe.fill.fore_color.rgb = ACCENT; stripe.line.fill.background()

tb = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(12), Inches(1.5))
tf = tb.text_frame
p = tf.paragraphs[0]
r = p.add_run(); r.text = "PathWise AI"
r.font.size = Pt(72); r.font.bold = True; r.font.color.rgb = WHITE

tb = s.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(12), Inches(0.7))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "An Intelligent, Vendor-Agnostic SD-WAN Management Platform"
r.font.size = Pt(24); r.font.color.rgb = TEAL; r.font.italic = True

tb = s.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(12), Inches(3))
tf = tb.text_frame
for line, sz, color, bold in [
    ("Team Pathfinders", 22, ACCENT, True),
    ("COSC 6370-001  |  Advanced Software Engineering", 16, WHITE, False),
    ("", 8, WHITE, False),
    ("Vineeth Reddy Kodakandla  ·  Meghana Nalluri", 14, WHITE, False),
    ("Bharadwaj Jakkula  ·  Sricharitha Katta", 14, WHITE, False),
    ("", 8, WHITE, False),
    ("April 16, 2026", 14, TEAL, False),
]:
    if line == "":
        p = tf.add_paragraph() if tf.paragraphs[0].runs else tf.paragraphs[0]
        p.space_after = Pt(2)
        continue
    if not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    r = p.add_run(); r.text = line
    r.font.size = Pt(sz); r.font.color.rgb = color; r.font.bold = bold

# =========================================================================
# SLIDE 2 — THE PROBLEM
# =========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_bar(s)
add_title(s, "The Problem", "The \"Switching Gap\" in Enterprise Networks")

add_text_block(s,
    "Modern SD-WAN solutions are REACTIVE — they failover only AFTER a link has already failed.",
    0.6, 1.55, 12.1, 0.7, size=16, bold=True, color=NAVY)

add_text_block(s,
    "During that window — between degradation and failover — packets drop. VoIP calls cut out. Video freezes. Financial transactions fail.",
    0.6, 2.25, 12.1, 0.9, size=14, color=GRAY)

# pain points cards
pains = [
    ("No Predictive Intelligence", "Reactive systems only detect failures after they occur."),
    ("Vendor Lock-In", "Proprietary appliances from Cisco/Juniper at premium cost."),
    ("CLI-Only Configuration", "Requires specialized network engineers."),
    ("Packet Loss on Failover", "Drops during hand-off disrupt mission-critical flows."),
]
for i, (t, b) in enumerate(pains):
    add_card(s, 0.6 + (i % 2) * 6.1, 3.4 + (i // 2) * 1.6, 5.9, 1.4, t, b, accent=RED)

add_footer(s, 2)

# =========================================================================
# SLIDE 3 — OUR SOLUTION
# =========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_bar(s)
add_title(s, "Our Solution", "Reactive  →  Predictive")

solutions = [
    ("30–60 sec Forecast", "LSTM neural networks predict WAN degradation before it impacts users."),
    ("Autonomous Steering", "Routes traffic via SDN controllers (OpenDaylight + ONOS)."),
    ("Hitless Handoff", "Zero packet loss, zero session drops during failover."),
    ("Vendor-Agnostic", "Runs on commodity x86-64 with any OpenFlow 1.3+ device."),
    ("Digital Twin Safety", "Every change validated in a sandbox before deployment."),
    ("Natural Language Policy", "No CLI. Type intent, system generates YANG/NETCONF."),
]
for i, (t, b) in enumerate(solutions):
    add_card(s, 0.6 + (i % 3) * 4.12, 1.75 + (i // 3) * 2.5, 3.95, 2.3, t, b, accent=TEAL)

add_footer(s, 3)

# =========================================================================
# SLIDE 4 — TARGET USERS
# =========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_bar(s)
add_title(s, "Target Users", "Who benefits from PathWise AI")

headers = ["User Segment", "Why They Need PathWise"]
rows = [
    ("SMEs", "Enterprise-grade reliability without enterprise budgets."),
    ("Managed Service Providers (MSPs)", "Centrally manage 100+ client sites from one dashboard."),
    ("Healthcare Facilities", "HIPAA-compliant, zero-downtime telemedicine + EHR."),
    ("Educational Institutions", "Reliable remote learning and campus-wide VoIP."),
    ("Retail Chains", "Continuous POS / payment uptime across all branches."),
]
add_table(s, headers, rows, 0.6, 1.75, 12.1, 4.5, font_size=14)
add_footer(s, 4)

# =========================================================================
# SLIDE 5 — ARCHITECTURE
# =========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_bar(s)
add_title(s, "System Architecture", "7-Service Microservices Platform (Dockerized)")

# Top layer - Frontend
def add_box(slide, left, top, w, h, text, fill, text_color=WHITE, size=11, bold=True):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = text_color; r.font.bold = bold

# Frontend row
add_box(s, 1.0, 1.7, 4.5, 0.7, "Dashboard (React + D3.js)", NAVY)
add_box(s, 7.8, 1.7, 4.5, 0.7, "IBN Interface (Natural Language)", NAVY)

# API layer
add_box(s, 2.5, 2.7, 8.3, 0.7, "Backend API  —  FastAPI + JWT + RBAC + WebSocket", TEAL)

# Services row
add_box(s, 0.5, 3.8, 3.0, 0.9, "Telemetry Engine\n(LSTM + Attention)", ACCENT, size=11)
add_box(s, 3.7, 3.8, 3.0, 0.9, "Traffic Steering\n(ODL + ONOS)", ACCENT, size=11)
add_box(s, 6.9, 3.8, 3.0, 0.9, "Digital Twin\n(Mininet + Batfish)", ACCENT, size=11)
add_box(s, 10.1, 3.8, 2.8, 0.9, "Redis Broker\n(Pub/Sub)", ACCENT, size=11)

# Data layer
add_box(s, 1.5, 5.1, 10.3, 0.8, "TimescaleDB  —  Telemetry  |  Health Scores  |  Audit Log  |  Users  |  Policies", NAVY, size=13)

# connector lines (simple)
from pptx.oxml.ns import qn
from lxml import etree

def add_line(slide, x1, y1, x2, y2, color=GRAY, weight=1.5):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(weight)

add_line(s, 3.25, 2.4, 4.5, 2.7)
add_line(s, 10.0, 2.4, 8.5, 2.7)
add_line(s, 6.65, 3.4, 2.0, 3.8)
add_line(s, 6.65, 3.4, 5.2, 3.8)
add_line(s, 6.65, 3.4, 8.4, 3.8)
add_line(s, 6.65, 3.4, 11.5, 3.8)
add_line(s, 2.0, 4.7, 6.65, 5.1)
add_line(s, 5.2, 4.7, 6.65, 5.1)

# legend / note
add_text_block(s,
    "Inter-service events flow over Redis pub/sub (telemetry, alerts, validation, steering, dashboard updates).",
    0.6, 6.2, 12.1, 0.6, size=12, color=GRAY, align=PP_ALIGN.CENTER)

add_footer(s, 5)

# =========================================================================
# SLIDE 6 — TECHNOLOGY STACK
# =========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_bar(s)
add_title(s, "Technology Stack", "Best-in-class open-source components")

headers = ["Layer", "Technology", "Purpose"]
rows = [
    ("ML / Prediction", "Python 3.11 + PyTorch 2.x", "LSTM with Bahdanau attention"),
    ("Time-Series DB", "TimescaleDB (PostgreSQL)", "Telemetry + health score storage"),
    ("Message Broker", "Redis 7.x Pub/Sub", "Real-time inter-service events"),
    ("Backend", "FastAPI + Uvicorn", "Async REST + WebSocket APIs"),
    ("SDN Control", "OpenDaylight + ONOS", "Northbound flow-table management"),
    ("Emulation", "Mininet (via WSL2)", "Digital twin + training data"),
    ("Validation", "Batfish", "Loop + firewall policy checks"),
    ("Frontend", "React 18 + TypeScript + D3.js v7", "SPA dashboard + visualizations"),
    ("Security", "JWT, bcrypt, TLS 1.3, AES-256", "Auth + encryption in transit/at rest"),
    ("Deployment", "Docker + Docker Compose", "Containerized microservices"),
]
add_table(s, headers, rows, 0.6, 1.7, 12.1, 5.3, font_size=12)
add_footer(s, 6)

# =========================================================================
# SLIDE 7 — LSTM PREDICTION ENGINE
# =========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_bar(s)
add_title(s, "The LSTM Prediction Engine", "PathWiseLSTM  —  The brain of the platform")

# left: architecture
add_text_block(s, "Architecture", 0.6, 1.7, 6, 0.4, size=18, bold=True, color=TEAL)
arch_items = [
    "Input:  60 timesteps @ 1 Hz (latency, jitter, packet loss)",
    "2 stacked LSTM layers, hidden size = 128",
    "Dropout = 0.2 between layers",
    "Bahdanau attention over sequence output",
    "Output:  predictions at t+30s AND t+60s",
    "Health score:  weighted combination → 0–100 scale",
]
add_bullets(s, arch_items, left=0.6, top=2.2, width=6, height=3, size=13, spacing=8)

# right: training
add_text_block(s, "Training", 6.9, 1.7, 6, 0.4, size=18, bold=True, color=ACCENT)
train_items = [
    "Optimizer:  Adam,  lr = 1e-3",
    "Loss:  MSE over 6 prediction targets",
    "Batch size:  256,  Epochs:  50 (early stop p=5)",
    "Dataset:  ≥ 10M synthetic Mininet points",
    "Train / Val / Test split:  70 / 15 / 15",
    "Target accuracy:  ≥ 90%  (Req-Qual-Perf-1)",
]
add_bullets(s, train_items, left=6.9, top=2.2, width=6, height=3, size=13, spacing=8)

# bottom highlight
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.1))
box.fill.solid(); box.fill.fore_color.rgb = NAVY
box.line.fill.background()
tf = box.text_frame
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Inference latency < 1 second per pass  —  fits inside the 1 Hz telemetry polling loop"
r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE

add_footer(s, 7)

# =========================================================================
# SLIDE 8 — HITLESS HANDOFF
# =========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_bar(s)
add_title(s, "Hitless Handoff Sequence", "Trigger → SDN flow update in < 50 milliseconds")

steps = [
    ("1", "LSTM predicts degradation; health score drops below threshold"),
    ("2", "Routing proposal sent to Digital Twin Sandbox"),
    ("3", "Sandbox validates: Mininet topology + Batfish loop/ACL check"),
    ("4", "PASSED → Traffic Steering engine activates"),
    ("5", "Pre-compute flow entries for alternative link"),
    ("6", "Atomic flow-table update via ODL/ONOS northbound API"),
    ("7", "Preserve TCP/VoIP session states throughout"),
    ("8", "Read-back confirmation + tamper-evident audit log entry"),
]
for i, (n, text) in enumerate(steps):
    row = i // 2
    col = i % 2
    left = 0.6 + col * 6.2
    top = 1.75 + row * 0.95
    # number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(0.65), Inches(0.65))
    circle.fill.solid(); circle.fill.fore_color.rgb = TEAL
    circle.line.fill.background()
    tf = circle.text_frame
    tf.margin_left=Inches(0); tf.margin_right=Inches(0)
    tf.margin_top=Inches(0.05); tf.margin_bottom=Inches(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = n
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE
    # step text
    add_text_block(s, text, left + 0.85, top + 0.08, 5.3, 0.6, size=12, color=NAVY)

# footer highlight
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.7))
box.fill.solid(); box.fill.fore_color.rgb = ACCENT
box.line.fill.background()
tf = box.text_frame; tf.margin_top = Inches(0.08)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "SLA:  < 50 ms end-to-end  (Req-Qual-Perf-2)  ·  Zero packet loss confirmed in integration tests"
r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE

add_footer(s, 8)

# =========================================================================
# SLIDE 9 — DIGITAL TWIN SANDBOX
# =========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_bar(s)
add_title(s, "Digital Twin Sandbox", "Validate every AI decision before production")

add_text_block(s,
    "Every routing change is auto-submitted. No exceptions. No direct-to-production deployments.",
    0.6, 1.65, 12.1, 0.5, size=14, bold=True, color=NAVY)

# Validation flow
steps = ["Mininet\nTopology", "Apply\nProposed Change", "Loop\nDetection (DFS)", "Batfish\nACL + Firewall", "PASSED?"]
for i, step in enumerate(steps):
    left = 0.6 + i * 2.5
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.4), Inches(2.2), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = TEAL if i < 4 else ACCENT
    box.line.fill.background()
    tf = box.text_frame; tf.margin_top = Inches(0.15); tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = step
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE
    if i < 4:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left + 2.2), Inches(2.8), Inches(0.3), Inches(0.4))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = GRAY
        arrow.line.fill.background()

# Two outcomes
pass_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.2), Inches(5.5), Inches(1.8))
pass_box.fill.solid(); pass_box.fill.fore_color.rgb = GREEN
pass_box.line.fill.background()
tf = pass_box.text_frame; tf.margin_left=Inches(0.2); tf.margin_top=Inches(0.2); tf.word_wrap=True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "✓  PASSED"
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.space_before = Pt(6)
r2 = p2.add_run(); r2.text = "Forward change payload to Traffic Steering engine for immediate deployment."
r2.font.size = Pt(12); r2.font.color.rgb = WHITE

fail_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(4.2), Inches(5.5), Inches(1.8))
fail_box.fill.solid(); fail_box.fill.fore_color.rgb = RED
fail_box.line.fill.background()
tf = fail_box.text_frame; tf.margin_left=Inches(0.2); tf.margin_top=Inches(0.2); tf.word_wrap=True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "✗  FAILED"
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.space_before = Pt(6)
r2 = p2.add_run(); r2.text = "Reject change. Log specific violation. Alert admin. Never touches production."
r2.font.size = Pt(12); r2.font.color.rgb = WHITE

# SLA banner
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.6))
box.fill.solid(); box.fill.fore_color.rgb = NAVY
box.line.fill.background()
tf = box.text_frame; tf.margin_top = Inches(0.05)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "SLA:  Full validation cycle < 5 seconds  (Req-Qual-Perf-3)"
r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE

add_footer(s, 9)

# =========================================================================
# SLIDE 10 — IBN (Intent-Based Networking)
# =========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_bar(s)
add_title(s, "Intent-Based Networking (IBN)", "Type your intent — no CLI required")

# Example input box
inp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.7), Inches(12.1), Inches(1.0))
inp.fill.solid(); inp.fill.fore_color.rgb = LIGHT
inp.line.color.rgb = TEAL; inp.line.width = Pt(2)
tf = inp.text_frame; tf.margin_left=Inches(0.3); tf.margin_top=Inches(0.15); tf.word_wrap=True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "USER INPUT:  "
r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = GRAY
r2 = p.add_run(); r2.text = "\"Prioritize VoIP traffic over video streaming on the Fiber link.\""
r2.font.size = Pt(16); r2.font.italic = True; r2.font.color.rgb = NAVY

# Flow
stages = [
    ("NLP Parser", "Regex + embedding similarity.\nNo external LLM dependency."),
    ("Intent Preview", "User reviews interpreted policy\nbefore Apply. Never deploys blindly."),
    ("YANG Translator", "ietf-diffserv-classifier\nietf-qos-policy modules."),
    ("SDN Controller", "NETCONF payload submitted\nto ODL/ONOS for deployment."),
]
for i, (t, b) in enumerate(stages):
    add_card(s, 0.6 + i * 3.12, 3.0, 2.95, 1.9, t, b, accent=ACCENT)

# Controlled vocabulary note
add_text_block(s, "Supported intent patterns", 0.6, 5.2, 12, 0.4, size=14, bold=True, color=TEAL)
patterns = [
    "prioritize <traffic> over <traffic>",
    "block <traffic> from <scope>",
    "limit <traffic> to <bandwidth>",
    "route <traffic> via <link_type>",
]
add_bullets(s, patterns, left=0.6, top=5.6, width=12, height=1.5, size=12, spacing=4)

add_footer(s, 10)

# =========================================================================
# SLIDE 11 — DASHBOARD UI
# =========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_bar(s)
add_title(s, "Dashboard UI", "Six key screens  —  React 18 + TypeScript + D3.js")

screens = [
    ("Login", "TLS 1.3 only · Generic error messages · 5-strike account lockout"),
    ("Telemetry Monitor", "Real-time D3 line graphs · Per-link filters · Time range selector"),
    ("Health Scoreboard", "Color-coded cards per link (Fiber/Sat/5G/Broadband) · Confidence + reasoning"),
    ("Policy Manager (IBN)", "Natural language input · Interpreted policy preview · Active policy list"),
    ("Audit Log", "Paginated · Filterable by event/date/actor · Tamper-evident checksums"),
    ("Reports", "PDF + CSV export of historical health scores, accuracy, steering events"),
]
for i, (t, b) in enumerate(screens):
    row = i // 3; col = i % 3
    add_card(s, 0.6 + col * 4.12, 1.75 + row * 2.5, 3.95, 2.3, t, b, accent=TEAL)

# footer highlight
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.75), Inches(12.1), Inches(0.3))
box.fill.solid(); box.fill.fore_color.rgb = ACCENT
box.line.fill.background()
tf = box.text_frame; tf.margin_top=Inches(0.02)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "UI response < 2 seconds under normal load  (Req-Qual-Perf-4)"
r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE

add_footer(s, 11)

# =========================================================================
# SLIDE 12 — SECURITY & COMPLIANCE
# =========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_bar(s)
add_title(s, "Security & Compliance", "Defense in depth at every layer")

headers = ["Requirement", "Implementation"]
rows = [
    ("Data in Transit", "TLS 1.3 enforced at nginx (lower versions rejected)"),
    ("Data at Rest", "AES-256 encryption for telemetry + credentials"),
    ("Passwords", "bcrypt one-way hash — never stored in plaintext"),
    ("Authentication", "JWT tokens, 60-minute expiry"),
    ("Authorization (RBAC)", "5 roles: NetAdmin, IT Manager, MSP, IT Staff, End User"),
    ("Account Lockout", "5 consecutive failed attempts triggers lock"),
    ("HIPAA Compliance", "Tamper-evident SHA-256 audit checksums"),
    ("Audit Completeness", "Every steering, policy, and auth event logged"),
]
add_table(s, headers, rows, 0.6, 1.75, 12.1, 5.2, font_size=13)
add_footer(s, 12)

# =========================================================================
# SLIDE 13 — DATA ARCHITECTURE
# =========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_bar(s)
add_title(s, "Data Architecture", "TimescaleDB + Redis Pub/Sub")

# Left: tables
add_text_block(s, "TimescaleDB Schema (5 tables)", 0.6, 1.65, 6, 0.4, size=16, bold=True, color=TEAL)
tables = [
    ("wan_telemetry", "Hypertable · 1-sec granularity (latency, jitter, loss)"),
    ("health_scores", "LSTM output + confidence per link"),
    ("audit_log", "Tamper-evident with SHA-256 row checksums"),
    ("users", "bcrypt credentials + RBAC role"),
    ("policies", "Natural language + YANG JSON config"),
]
for i, (name, desc) in enumerate(tables):
    top = 2.2 + i * 0.85
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(top), Inches(6), Inches(0.75))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = TEAL; box.line.width = Pt(1)
    tf = box.text_frame; tf.margin_left=Inches(0.15); tf.margin_top=Inches(0.05); tf.word_wrap=True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = name
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = desc
    r2.font.size = Pt(10); r2.font.color.rgb = GRAY

# Right: Redis channels
add_text_block(s, "Redis Pub/Sub Channels", 7.0, 1.65, 6, 0.4, size=16, bold=True, color=ACCENT)
channels = [
    "pathwise:telemetry:{link_id}",
    "pathwise:alerts:{site_id}",
    "pathwise:validation:request",
    "pathwise:validation:result",
    "pathwise:steering:trigger",
    "pathwise:dashboard:updates",
]
for i, ch in enumerate(channels):
    top = 2.2 + i * 0.7
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(top), Inches(5.8), Inches(0.6))
    box.fill.solid(); box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    tf = box.text_frame; tf.margin_left=Inches(0.2); tf.margin_top=Inches(0.08)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = ch
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = TEAL
    r.font.name = "Consolas"

add_footer(s, 13)

# =========================================================================
# SLIDE 14 — QUALITY TARGETS
# =========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_bar(s)
add_title(s, "Quality Targets", "Hard constraints  —  every target is test-verified")

headers = ["Metric", "Target", "SRS ID"]
rows = [
    ("LSTM prediction accuracy", "≥ 90% (MSE on test set)", "Req-Qual-Perf-1"),
    ("End-to-end traffic steering", "< 50 ms", "Req-Qual-Perf-2"),
    ("Digital Twin validation cycle", "< 5 seconds", "Req-Qual-Perf-3"),
    ("Dashboard UI response", "< 2 seconds under load", "Req-Qual-Perf-4"),
    ("Concurrent sites supported", "≥ 100 (no degradation)", "Req-Qual-Scal-1"),
    ("Platform availability", "≥ 99.9% annually", "Req-Qual-Rel-1"),
    ("Data in transit", "TLS 1.3+ on all connections", "Req-Qual-Sec-1"),
    ("Data at rest", "AES-256 encryption", "Req-Qual-Sec-2"),
    ("Database backups", "Every 24h, geo-separated", "Req-Qual-Rel-3"),
]
add_table(s, headers, rows, 0.6, 1.75, 12.1, 5.2, font_size=12)
add_footer(s, 14)

# =========================================================================
# SLIDE 15 — TESTING STRATEGY
# =========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_bar(s)
add_title(s, "Testing Strategy", "22 mandatory integration test cases + ≥ 80% unit coverage")

tests = [
    ("Telemetry Ingestion", "≥ 1 row/sec into TimescaleDB"),
    ("LSTM Accuracy", "≥ 90% on held-out test set"),
    ("Inference Speed", "< 1 sec per LSTM pass"),
    ("Hitless Handoff", "< 50 ms, zero packet loss"),
    ("Session Preservation", "No TCP session dropped on handoff"),
    ("Sandbox Cycle", "Full validation < 5 seconds"),
    ("Loop Detection", "Correctly rejects loop-introducing changes"),
    ("IBN NLP Accuracy", "> 90% intent parsing on common commands"),
    ("YANG Translation", "Payload accepted by ODL + ONOS"),
    ("RBAC Enforcement", "Each role limited to permitted routes"),
    ("100-Site Scalability", "No degradation at load"),
    ("TLS 1.3 Enforcement", "Lower versions rejected"),
]
for i, (t, b) in enumerate(tests):
    row = i // 3; col = i % 3
    add_card(s, 0.6 + col * 4.12, 1.75 + row * 1.28, 3.95, 1.15, t, b, accent=TEAL)

add_footer(s, 15)

# =========================================================================
# SLIDE 16 — BUILD ORDER
# =========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_bar(s)
add_title(s, "Build Order", "10-stage sequence respecting service dependencies")

stages = [
    ("1", "Infrastructure", "TimescaleDB · Redis · nginx TLS"),
    ("2", "Backend API", "Auth · RBAC · REST routes · WebSocket"),
    ("3", "Telemetry Engine", "Ingestion · LSTM · scoring · alerts"),
    ("4", "Digital Twin", "Mininet builder · Batfish · sandbox API"),
    ("5", "Traffic Steering", "ODL/ONOS clients · hitless handoff"),
    ("6", "IBN Interface", "NLP engine · YANG translator · policy mgr"),
    ("7", "Dashboard", "React SPA · all 6 screens"),
    ("8", "ML Pipeline", "Data generation · training · evaluation"),
    ("9", "Integration Tests", "22 end-to-end test cases"),
    ("10", "Full Bring-Up", "Docker Compose · end-to-end demo"),
]
for i, (n, t, b) in enumerate(stages):
    row = i // 2; col = i % 2
    left = 0.6 + col * 6.2
    top = 1.7 + row * 1.05
    # num circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(0.7), Inches(0.7))
    circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    tf = circle.text_frame; tf.margin_top = Inches(0.07)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = n
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE
    # text
    tb = s.shapes.add_textbox(Inches(left + 0.85), Inches(top + 0.02), Inches(5.2), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = t
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = b
    r2.font.size = Pt(10); r2.font.color.rgb = GRAY

add_footer(s, 16)

# =========================================================================
# SLIDE 17 — CHALLENGES & MITIGATIONS
# =========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_bar(s)
add_title(s, "Challenges & Mitigations", "Anticipating the hard parts")

headers = ["Challenge", "Mitigation Strategy"]
rows = [
    ("50 ms hitless handoff is extremely tight", "Pre-compute flow entries · atomic SDN update · session-state preservation"),
    ("Running Mininet on Windows host", "Execute via WSL2 subprocess wrapper from containerized service"),
    ("Achieving ≥ 90% LSTM accuracy", "10M synthetic training points · attention mechanism · hyperparameter tuning"),
    ("Ambiguous natural language intents", "Preview before Apply · rephrasing suggestions · never guess or auto-deploy"),
    ("SDN controller API failures", "Retry with exponential backoff (3 attempts) · admin alert on failure"),
    ("Tamper-evident audit logging", "SHA-256 checksum per audit_log row · immutable append-only storage"),
    ("Service restart / reconnection", "Docker restart:always · health checks · Redis reconnection logic"),
]
add_table(s, headers, rows, 0.6, 1.75, 12.1, 5.0, font_size=12)
add_footer(s, 17)

# =========================================================================
# SLIDE 18 — DEMO + Q&A
# =========================================================================
s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.0), prs.slide_width, Inches(0.08))
stripe.fill.solid(); stripe.fill.fore_color.rgb = ACCENT; stripe.line.fill.background()

tb = s.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(12), Inches(1.2))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Live Demo  &  Q & A"
r.font.size = Pt(56); r.font.bold = True; r.font.color.rgb = WHITE

tb = s.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(12), Inches(0.6))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Seeing PathWise AI in action"
r.font.size = Pt(22); r.font.color.rgb = TEAL; r.font.italic = True

# demo flow
tb = s.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(12), Inches(0.5))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "DEMO FLOW"
r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = ACCENT

demo_steps = [
    "1.  Login → Health Scoreboard displays live Fiber / Satellite / 5G / Broadband",
    "2.  Inject simulated degradation on Fiber link",
    "3.  LSTM predicts → health score drops below threshold",
    "4.  Digital Twin Sandbox validates → PASSED",
    "5.  Traffic steering reroutes → zero packet loss confirmed",
    "6.  Audit log entry appears with confidence score + reasoning",
]
tb = s.shapes.add_textbox(Inches(0.8), Inches(3.9), Inches(12), Inches(3))
tf = tb.text_frame
for i, step in enumerate(demo_steps):
    if i == 0: p = tf.paragraphs[0]
    else: p = tf.add_paragraph()
    p.space_after = Pt(5)
    r = p.add_run(); r.text = step
    r.font.size = Pt(14); r.font.color.rgb = WHITE

# thanks
tb = s.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(12), Inches(0.5))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Thank you  ·  Questions?"
r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = ACCENT


# Save
out = r"C:\Users\vinee\Desktop\PATHWISEAI\PathWise_AI_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
